from typing import List, Union, Optional
from pydantic import BaseModel, Field
from indexify_extractor_sdk import Content, Extractor, Feature
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests
import tempfile

class PPTExtractorConfig(BaseModel):
    output_types: List[str] = Field(default_factory=lambda: ["text", "table"])

class PPTExtractor(Extractor):
    name = "tensorlake/ppt"
    description = "An extractor that let's you extract information from presentations."
    system_dependencies = []
    input_mime_types = ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]

    def __init__(self):
        super(PPTExtractor, self).__init__()

    def extract(self, content: Content, params: PPTExtractorConfig) -> List[Union[Feature, Content]]:
        contents = []
        prs = None
        with tempfile.NamedTemporaryFile(delete=True, suffix=".pptx") as inputtmpfile:
            inputtmpfile.write(content.data)
            inputtmpfile.flush()
            prs = Presentation(inputtmpfile.name)
        
            # Iterate through slides
        for slide_idx, slide in enumerate(prs.slides):
            # Extract text
            text_output = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_output.append(paragraph.text)
                    contents.append(Content.from_text("\n".join(text_output), features=[Feature.metadata({"page": slide_idx})]))

            # Extract images
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    contents.append(Content(content_type="image/png", data=image.blob, features=[Feature.metadata({"page": slide_idx})]))

            # Extract tables
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)

                    # Save table to a file
                    contents.append(Content.from_json(table_data, features=[Feature.metadata({"page": slide_idx})]))

        return contents

    def sample_input(self) -> Content:
        req = requests.get("https://pub-157277cc11d64fb1a11f71cc52c688eb.r2.dev/figures.pptx")
        ppt_data = req.content
        return Content(content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", data=ppt_data)

if __name__ == "__main__":
    filepath = "test.pptx"
    with open(filepath, 'rb') as f:
        ppt_data = f.read()
    data = Content(content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", data=ppt_data)
    extractor = PPTExtractor()
    params = PPTExtractorConfig(output_types=["text", "table"])
    results = extractor.extract(data, params)
    print(results)